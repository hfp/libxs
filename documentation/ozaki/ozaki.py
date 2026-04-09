#!/usr/bin/env python3
###############################################################################
# Copyright (c) Intel Corporation - All rights reserved.                      #
#                                                                             #
# For information on the license, see the LICENSE file.                       #
# SPDX-License-Identifier: BSD-3-Clause                                       #
###############################################################################
import os
import re
import subprocess

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# Slide dimensions (16:9 widescreen, 1920x1080 equivalent)
SLIDE_WIDTH = Inches(40 / 3)  # 13+1/3"
SLIDE_HEIGHT = Inches(7.5)

# Fonts
CODE_FONT = "Courier New"
CODE_SIZE = Pt(20)
TABLE_SIZE = Pt(12)

# Layout
MARGIN = Inches(0.5)
TITLE_TOP = Inches(0.5)
TITLE_BOX = Inches(1.25)
BODY_TOP = TITLE_TOP + TITLE_BOX + Inches(0.1)
BODY_HEIGHT = SLIDE_HEIGHT - BODY_TOP - MARGIN
CONTENT_WIDTH = SLIDE_WIDTH - 2 * MARGIN
TABLE_ROW_HEIGHT = Inches(0.5)
LINE_HEIGHT = Inches(0.25)
BLOCK_SPACING = Inches(0.2)


# Inline markdown parsing
#
def parse_inline(text):
    """Parse **bold** and `code` into [(text, style), ...]."""
    result = []
    pos = 0
    for m in re.finditer(r"\*\*(.+?)\*\*|`(.+?)`", text):
        if m.start() > pos:
            result.append((text[pos : m.start()], "normal"))
        if m.group(1) is not None:
            result.append((m.group(1), "bold"))
        else:
            result.append((m.group(2), "code"))
        pos = m.end()
    if pos < len(text):
        result.append((text[pos:], "normal"))
    return result or [("", "normal")]


def add_runs(para, text, font_size=None, code=False):
    """Add inline-formatted runs to a paragraph."""
    if code:
        run = para.add_run()
        run.text = text
        run.font.name = CODE_FONT
        if font_size:
            run.font.size = font_size
        return
    for content, style in parse_inline(text):
        run = para.add_run()
        run.text = content
        if font_size:
            run.font.size = font_size
        if style == "bold":
            run.font.bold = True
        elif style == "code":
            run.font.name = CODE_FONT


def no_bullet(para):
    """Suppress bullet formatting on a placeholder paragraph."""
    pPr = para._p.get_or_add_pPr()
    pPr.set("indent", "0")
    pPr.set("marL", "0")
    for tag in ("buNone", "buChar", "buAutoNum"):
        for el in pPr.findall(qn(f"a:{tag}")):
            pPr.remove(el)
    etree.SubElement(pPr, qn("a:buNone"))


# Markdown parsing
#
def parse_markdown(path):
    """Split markdown on --- separators and parse each slide."""
    with open(path) as f:
        raw = f.read()
    return [
        s
        for s in (_parse_slide(part.strip()) for part in re.split(r"\n---\n", raw))
        if s
    ]


def _skip_blank(lines, i):
    while i < len(lines) and not lines[i].strip():
        i += 1
    return i


def _parse_slide(text):
    if not text:
        return None
    lines = text.split("\n")
    slide = {"title": "", "subtitle": "", "is_title": False, "blocks": []}
    i = _skip_blank(lines, 0)

    # Detect title slide (# heading, not ##)
    if i < len(lines) and re.match(r"^# (?!#)", lines[i]):
        slide["title"] = lines[i][2:].strip()
        slide["is_title"] = True
        i = _skip_blank(lines, i + 1)
        if i < len(lines) and lines[i].startswith("## "):
            slide["subtitle"] = lines[i][3:].strip()
            i += 1
    elif i < len(lines) and lines[i].startswith("## "):
        slide["title"] = lines[i][3:].strip()
        i += 1

    # Parse body into blocks
    cur = None
    in_code = False
    while i < len(lines):
        line = lines[i]

        # Code fence toggle
        if line.strip().startswith("```"):
            if in_code:
                slide["blocks"].append(cur)
                cur, in_code = None, False
            else:
                if cur:
                    slide["blocks"].append(cur)
                    cur = None
                cur = {"type": "code", "lines": []}
                in_code = True
            i += 1
            continue

        if in_code:
            cur["lines"].append(line)
            i += 1
            continue

        # Table row
        if line.strip().startswith("|") and "|" in line.strip()[1:]:
            if cur and cur["type"] != "table":
                slide["blocks"].append(cur)
                cur = None
            if not cur:
                cur = {"type": "table", "rows": []}
            if not re.match(r"^\s*\|[\s\-:|]+\|\s*$", line):
                cells = [c.strip() for c in line.strip().strip("|").split("|")]
                cur["rows"].append(cells)
            i += 1
            continue

        # Sub-heading (### within a slide)
        if line.startswith("### "):
            if cur:
                slide["blocks"].append(cur)
                cur = None
            slide["blocks"].append({"type": "heading", "text": line[4:].strip()})
            i += 1
            continue

        # Bullet item
        if line.strip().startswith("- "):
            if cur and cur["type"] != "bullets":
                slide["blocks"].append(cur)
                cur = None
            if not cur:
                cur = {"type": "bullets", "items": []}
            cur["items"].append(line.strip()[2:])
            i += 1
            continue

        # Empty line flushes current block
        if not line.strip():
            if cur:
                slide["blocks"].append(cur)
                cur = None
            i += 1
            continue

        # Regular text (including numbered lists)
        if cur and cur["type"] == "text":
            cur["lines"].append(line.strip())
        else:
            if cur:
                slide["blocks"].append(cur)
            cur = {"type": "text", "lines": [line.strip()]}
        i += 1

    if cur:
        slide["blocks"].append(cur)
    return slide


# Presentation building
#
def build_presentation(slides):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    # Resize placeholders to match slide dimensions
    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            idx = ph.placeholder_format.idx
            if idx <= 1:
                ph.left = MARGIN
                ph.width = CONTENT_WIDTH
        # Adjust vertical positions for non-title-slide layouts
        if layout != prs.slide_layouts[0]:
            for ph in layout.placeholders:
                idx = ph.placeholder_format.idx
                if idx == 0:
                    ph.top = TITLE_TOP
                    ph.height = TITLE_BOX
                elif idx == 1:
                    ph.top = BODY_TOP
                    ph.height = BODY_HEIGHT
    for s in slides:
        if s["is_title"]:
            _title_slide(prs, s)
        elif any(b["type"] == "table" for b in s["blocks"]):
            _table_slide(prs, s)
        else:
            _content_slide(prs, s)
    return prs


def _title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.placeholders[0].text = data["title"]
    tf = slide.placeholders[1].text_frame
    if data["subtitle"]:
        tf.text = data["subtitle"]
    for block in data["blocks"]:
        if block["type"] == "text":
            for line in block["lines"]:
                p = tf.add_paragraph()
                add_runs(p, line)


def _content_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.placeholders[0].text = data["title"]
    tf = slide.placeholders[1].text_frame
    tf.clear()
    _render_blocks(tf, data["blocks"], placeholder=True)


def _table_slide(prs, data):
    """Slide containing a table: Title Only layout with manual shapes."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.placeholders[0].text = data["title"]
    left = MARGIN
    top = BODY_TOP
    width = CONTENT_WIDTH

    for block in data["blocks"]:
        if block["type"] == "table":
            top = _add_table(slide, block["rows"], left, top, width)
        else:
            n = 1
            if block["type"] in ("code", "bullets"):
                n = len(block.get("lines", block.get("items", [])))
            elif block["type"] == "text":
                n = len(block["lines"])
            height = LINE_HEIGHT * max(n, 1)
            txbox = slide.shapes.add_textbox(left, top, width, height)
            tf = txbox.text_frame
            tf.word_wrap = True
            _render_blocks(tf, [block], placeholder=False)
            top += height + BLOCK_SPACING


def _render_blocks(tf, blocks, placeholder=False):
    """Render block list into a text frame."""
    first = True
    for block in blocks:
        btype = block["type"]

        if btype == "text":
            for line in block["lines"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                add_runs(p, line)
                if placeholder:
                    no_bullet(p)

        elif btype == "heading":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = block["text"]
            run.font.bold = True
            if placeholder:
                no_bullet(p)

        elif btype == "bullets":
            for item in block["items"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                if not placeholder:
                    # Textboxes don't have automatic bullets
                    item = "\u2022 " + item
                add_runs(p, item)

        elif btype == "code":
            for line in block["lines"]:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                add_runs(p, line, font_size=CODE_SIZE, code=True)
                if placeholder:
                    no_bullet(p)
                p.space_before = Pt(0)
                p.space_after = Pt(0)


def _add_table(slide, rows, left, top, width):
    """Add a table shape and return the new vertical position."""
    if not rows:
        return top
    n_rows, n_cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(
        n_rows, n_cols, left, top, width, TABLE_ROW_HEIGHT * n_rows
    )
    tbl = shape.table
    for r, row_data in enumerate(rows):
        for c, cell_text in enumerate(row_data):
            if c < n_cols:
                p = tbl.cell(r, c).text_frame.paragraphs[0]
                add_runs(p, cell_text, font_size=TABLE_SIZE)
                if r == 0:
                    for run in p.runs:
                        run.font.bold = True
    return top + TABLE_ROW_HEIGHT * n_rows + BLOCK_SPACING


# Post-processing
#
def autofit(pptx_path):
    """Add normAutofit to every text frame that lacks an autofit setting."""
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                body_pr = shape.text_frame._txBody.find(qn("a:bodyPr"))
                if body_pr is None:
                    continue
                if (
                    body_pr.find(qn("a:normAutofit")) is not None
                    or body_pr.find(qn("a:spAutoFit")) is not None
                    or body_pr.find(qn("a:noAutofit")) is not None
                ):
                    continue
                etree.SubElement(body_pr, qn("a:normAutofit"))
    prs.save(pptx_path)


def resave(pptx_path):
    """Toggle auto-fit on each shape via PowerPoint COM to force reflow."""
    winpath = subprocess.check_output(
        ["wslpath", "-w", os.path.abspath(pptx_path)],
        text=True,
    ).strip()
    # ppAutoSizeNone=0, ppAutoSizeTextToFitShape=2
    subprocess.check_call(
        [
            "powershell.exe",
            "-NoProfile",
            "-Command",
            (
                f"$pp = New-Object -ComObject PowerPoint.Application;"
                f'$pres = $pp.Presentations.Open("{winpath}");'
                f"foreach ($slide in $pres.Slides) {{"
                f"  foreach ($shape in $slide.Shapes) {{"
                f"    if ($shape.HasTextFrame) {{"
                f"      $shape.TextFrame2.AutoSize = 0;"
                f"      $shape.TextFrame2.AutoSize = 2"
                f"    }}"
                f"  }}"
                f"}}"
                f"$pres.Save();"
                f"$pres.Close();"
                f"$pp.Quit();"
                f"[System.Runtime.InteropServices.Marshal]"
                f"::ReleaseComObject($pp) | Out-Null"
            ),
        ]
    )


# Main
#
def main():
    here = os.path.dirname(os.path.abspath(__file__))
    output = os.path.join(here, "ozaki.pptx")
    source = os.path.join(here, "index.md")

    slides = parse_markdown(source)
    prs = build_presentation(slides)
    prs.save(output)
    autofit(output)
    resave(output)


if __name__ == "__main__":
    main()
